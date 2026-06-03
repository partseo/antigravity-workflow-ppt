import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ==========================================
# 테마 색상 및 스타일 정의 (레드 / 네이비 / 화이트)
# ==========================================
COLOR_PRIMARY_RED = RGBColor(219, 17, 36)      # 크림슨 레드 (포인트 & 강조)
COLOR_SECONDARY_NAVY = RGBColor(29, 37, 45)    # 다크 네이비 (메인 데코 및 텍스트)
COLOR_LIGHT_BG = RGBColor(245, 245, 245)       # 연한 그레이 (배경 및 카드 테두리/채우기)
COLOR_WHITE = RGBColor(255, 255, 255)          # 순수 흰색
COLOR_DARK_TEXT = RGBColor(60, 60, 60)         # 어두운 텍스트
COLOR_MUTED_TEXT = RGBColor(128, 128, 128)     # 뮤트된 텍스트

FONT_TITLE = "Malgun Gothic"
FONT_BODY = "Malgun Gothic"

# ==========================================
# 헬퍼 함수: 도형 스타일링
# ==========================================
def style_shape(shape, fill_color, line_color=None):
    """도형의 채우기 및 외곽선 색상을 지정하는 헬퍼 함수"""
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()

def set_text_formatting(paragraph, text, size_pt, color, bold=False, font_name=FONT_BODY, alignment=PP_ALIGN.LEFT):
    """텍스트 단락의 상세 포맷을 정의하는 헬퍼 함수"""
    paragraph.text = text
    paragraph.font.name = font_name
    paragraph.font.size = Pt(size_pt)
    paragraph.font.color.rgb = color
    paragraph.font.bold = bold
    paragraph.alignment = alignment

# ==========================================
# 헬퍼 함수: 공통 슬라이드 템플릿 데코레이션
# ==========================================
def add_common_slide_decorations(prs, slide, category_text, title_text):
    """모든 일반 본문 슬라이드에 적용되는 공통 헤더 및 모서리 데코레이션"""
    # 1. 좌상단 인덱스 뱃지 (크림슨 레드)
    badge = slide.shapes.add_shape(
        MSO_SHAPE.PARALLELOGRAM,
        Inches(0.5), Inches(0.4), Inches(1.8), Inches(0.35)
    )
    style_shape(badge, COLOR_PRIMARY_RED)
    # 평행사변형 각도 미세 조정 (기본값도 훌륭함)
    badge_tf = badge.text_frame
    badge_tf.word_wrap = False
    badge_tf.margin_left = Inches(0.1)
    badge_tf.margin_top = Inches(0.02)
    p0 = badge_tf.paragraphs[0]
    set_text_formatting(p0, category_text, 12, COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # 2. 메인 타이틀 텍스트
    title_box = slide.shapes.add_textbox(Inches(2.5), Inches(0.35), Inches(10.3), Inches(0.5))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_tf.margin_left = Inches(0)
    title_tf.margin_top = Inches(0)
    p_title = title_tf.paragraphs[0]
    set_text_formatting(p_title, title_text, 26, COLOR_SECONDARY_NAVY, bold=True, font_name=FONT_TITLE)
    
    # 3. 상단 얇은 장식선 (다크 네이비)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(0.9), Inches(12.33), Inches(0.03)
    )
    style_shape(line, COLOR_SECONDARY_NAVY)
    
    # 4. 우하단 모서리 장식 (크림슨 레드 삼각형)
    corner = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE,
        Inches(12.833), Inches(7.0), Inches(0.5), Inches(0.5)
    )
    style_shape(corner, COLOR_PRIMARY_RED)
    corner.rotation = 180 # 180도 회전시켜 우하단 모서리에 딱 맞춤

# ==========================================
# 12개 슬라이드 개별 생성 함수
# ==========================================

def build_slide_01_cover(prs):
    """Slide 1: 표지 (Cover)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 배경 도형 데코레이션 (템플릿 스타일 구현)
    # 좌상단 네이비 삼각형
    tri1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(-0.5), Inches(-0.5), Inches(4.5), Inches(4.5))
    style_shape(tri1, COLOR_SECONDARY_NAVY)
    tri1.rotation = 90
    
    # 좌상단 레드 삼각형
    tri2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(-0.2), Inches(-0.2), Inches(3.2), Inches(3.2))
    style_shape(tri2, COLOR_PRIMARY_RED)
    tri2.rotation = 90

    # 우하단 네이비 삼각형
    tri3 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(9.333), Inches(3.5), Inches(4.5), Inches(4.5))
    style_shape(tri3, COLOR_SECONDARY_NAVY)
    tri3.rotation = 270

    # 우하단 레드 삼각형
    tri4 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(10.333), Inches(4.5), Inches(3.2), Inches(3.2))
    style_shape(tri4, COLOR_PRIMARY_RED)
    tri4.rotation = 270

    # 중앙 타이틀 카드 백그라운드 (테두리 있는 라이트 그레이 카드)
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.0), Inches(1.8), Inches(9.33), Inches(3.8))
    style_shape(card, COLOR_WHITE, COLOR_LIGHT_BG)
    
    # 내부 테두리 데코용 라인 (상하 붉은색/네이비 얇은 라인)
    line_t = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.3), Inches(2.1), Inches(8.73), Inches(0.04))
    style_shape(line_t, COLOR_PRIMARY_RED)
    line_b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.3), Inches(5.2), Inches(8.73), Inches(0.04))
    style_shape(line_b, COLOR_SECONDARY_NAVY)

    # 타이틀 텍스트 박스
    title_box = slide.shapes.add_textbox(Inches(2.5), Inches(2.4), Inches(8.33), Inches(2.6))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    set_text_formatting(p0, "라면 자판기 기반 16:9 LCD 광고 플랫폼", 36, COLOR_SECONDARY_NAVY, bold=True, font_name=FONT_TITLE, alignment=PP_ALIGN.CENTER)
    
    p_space = tf.add_paragraph()
    p_space.font.size = Pt(14)
    
    p1 = tf.add_paragraph()
    set_text_formatting(p1, "식품 자동판매 + 생활공간 광고 네트워크 사업전략", 20, COLOR_PRIMARY_RED, bold=True, font_name=FONT_TITLE, alignment=PP_ALIGN.CENTER)
    
    p2 = tf.add_paragraph()
    set_text_formatting(p2, "\n투자자 IR 피치덱", 14, COLOR_MUTED_TEXT, bold=False, font_name=FONT_BODY, alignment=PP_ALIGN.CENTER)


def build_slide_02_toc(prs):
    """Slide 2: 목차 (Table of Contents)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 헤더 데코레이션
    badge = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, Inches(0.5), Inches(0.4), Inches(1.8), Inches(0.35))
    style_shape(badge, COLOR_SECONDARY_NAVY)
    badge_tf = badge.text_frame
    p0 = badge_tf.paragraphs[0]
    set_text_formatting(p0, "CONTENTS", 12, COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    title_box = slide.shapes.add_textbox(Inches(2.5), Inches(0.35), Inches(10.3), Inches(0.5))
    p_title = title_box.text_frame.paragraphs[0]
    set_text_formatting(p_title, "목차 및 주요 아젠다", 26, COLOR_SECONDARY_NAVY, bold=True, font_name=FONT_TITLE)
    
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.9), Inches(12.33), Inches(0.03))
    style_shape(line, COLOR_PRIMARY_RED)

    # 목차 4대 영역 구성 카드 (그리드 레이아웃)
    sections = [
        {"num": "01", "title": "시장 문제 및 해결방안", "desc": "소상공인 광고의 한계 극복 및\n조리 대기시간(2~3분) 매체화"},
        {"num": "02", "title": "제품 및 화면 구성 전략", "desc": "16:9 LCD 탑재 라면 자판기 및\n6단계 광고/콘텐츠 편성 계획"},
        {"num": "03", "title": "수익 모델 및 재무 시뮬레이션", "desc": "장비 마진 + 라면 마진 + 광고 배분\n1대 및 100대 기준 수익 검증"},
        {"num": "04", "title": "시장 기회 및 성장 로드맵", "desc": "대학, PC방 등 공간 확장 계획과\n전국망 광고 플랫폼으로의 비전"}
    ]
    
    for i, sec in enumerate(sections):
        col = i % 2
        row = i // 2
        
        left = Inches(1.0) if col == 0 else Inches(7.0)
        top = Inches(1.8) if row == 0 else Inches(4.3)
        
        # 카드 배경
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.3), Inches(2.1))
        style_shape(card, COLOR_WHITE, COLOR_LIGHT_BG)
        
        # 숫자 배지
        num_badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.3), top + Inches(0.3), Inches(0.8), Inches(0.8))
        style_shape(num_badge, COLOR_PRIMARY_RED)
        num_p = num_badge.text_frame.paragraphs[0]
        set_text_formatting(num_p, sec["num"], 16, COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        
        # 텍스트 박스
        text_box = slide.shapes.add_textbox(left + Inches(1.3), top + Inches(0.2), Inches(3.8), Inches(1.7))
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.1)
        
        p_t = tf.paragraphs[0]
        set_text_formatting(p_t, sec["title"], 18, COLOR_SECONDARY_NAVY, bold=True, font_name=FONT_TITLE)
        
        p_d = tf.add_paragraph()
        set_text_formatting(p_d, "\n" + sec["desc"], 12, COLOR_DARK_TEXT)


def build_slide_03_problem(prs):
    """Slide 3: Problem (시장 문제점)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_common_slide_decorations(prs, slide, "01. PROBLEM", "소상공인 광고 한계 및 자판기 수익성 정체")
    
    # 3단 세로형 카드 레이아웃
    problems = [
        {"title": "소상공인 광고의 높은 장벽", "point": "비싸고 효과 측정이 모호함", "desc": "기존 전단지, 버스, 지하철 광고 등은 초기 비용 부담이 크고, 타겟팅이나 노출 빈도 등의 직접적인 효과 측정이 불가능하여 소상공인 마케팅에 한계가 있음."},
        {"title": "자판기 운영자의 한정된 수익 구조", "point": "라면 판매 외 추가 수익 부재", "desc": "기존 자판기 운영은 식자재 마진(라면 등)에만 100% 의존하고 있어, 물가 상승이나 경쟁 심화 시 안정적인 추가 마진을 확보할 수 있는 대안 모델이 부재함."},
        {"title": "대기시간 광고 매체의 블루오션", "point": "조리 2~3분간의 강렬한 주목도", "desc": "소비자가 자판기 앞에서 음식이 조리되기를 기다리는 2~3분의 대기 시간은 매우 긴 체감 시간이지만, 시선을 끌어당기는 전용 광고 매체가 아직 대중화되지 않은 틈새 시장임."}
    ]
    
    for i, prob in enumerate(problems):
        left = Inches(0.5) + i * Inches(4.2)
        top = Inches(1.8)
        width = Inches(3.9)
        height = Inches(4.8)
        
        # 카드 배경
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        style_shape(card, COLOR_WHITE, COLOR_LIGHT_BG)
        
        # 카드의 탑 라인 (크림슨 레드)
        top_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.15))
        style_shape(top_line, COLOR_PRIMARY_RED)
        
        # 텍스트 내용
        text_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.3), width - Inches(0.4), height - Inches(0.5))
        tf = text_box.text_frame
        tf.word_wrap = True
        
        # 아이콘 대체 문자 및 제목
        p_num = tf.paragraphs[0]
        set_text_formatting(p_num, f"0{i+1}", 24, COLOR_PRIMARY_RED, bold=True)
        
        p_title = tf.add_paragraph()
        set_text_formatting(p_title, "\n" + prob["title"], 18, COLOR_SECONDARY_NAVY, bold=True, font_name=FONT_TITLE)
        
        p_point = tf.add_paragraph()
        set_text_formatting(p_point, "👉 " + prob["point"], 14, COLOR_PRIMARY_RED, bold=True)
        
        p_desc = tf.add_paragraph()
        set_text_formatting(p_desc, "\n" + prob["desc"], 12, COLOR_DARK_TEXT)


def build_slide_04_solution(prs):
    """Slide 4: Solution (해결 방안)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_common_slide_decorations(prs, slide, "01. SOLUTION", "조리 대기시간(2~3분) 매체화 및 3대 핵심 수익 모델 결합")
    
    # 좌측: 핵심 해결 방안 설명 (2단 구성)
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(6.0), Inches(4.8))
    ltf = left_box.text_frame
    ltf.word_wrap = True
    
    p0 = ltf.paragraphs[0]
    set_text_formatting(p0, "🎯 조리 대기시간 2~3분을 '광고 미디어'로 전환", 18, COLOR_PRIMARY_RED, bold=True)
    p0_sub = ltf.add_paragraph()
    set_text_formatting(p0_sub, "라면이 조리되는 동안 이용자의 시선은 자연스럽게 기기 화면에 초집중됩니다. 이 골든타임 동안 16:9 LCD 화면을 통해 타겟화된 광고를 100% 전달합니다.", 13, COLOR_DARK_TEXT)
    
    p1 = ltf.add_paragraph()
    set_text_formatting(p1, "\n📈 하이브리드 비즈니스 모델로 수익 극대화", 18, COLOR_SECONDARY_NAVY, bold=True)
    p1_sub = ltf.add_paragraph()
    set_text_formatting(p1_sub, "장비 판매로 빠른 현금 흐름을 만들고, 일상 소비재인 라면 판매 마진으로 안정적인 기반 매출을 내며, 매력적인 매체로서의 LCD 광고 플랫폼 반복 수익(Recurring Revenue)을 창출합니다.", 13, COLOR_DARK_TEXT)
    
    p2 = ltf.add_paragraph()
    set_text_formatting(p2, "\n🤝 본사·운영자·장소제공자 상생 생태계 구축", 18, COLOR_SECONDARY_NAVY, bold=True)
    p2_sub = ltf.add_paragraph()
    set_text_formatting(p2_sub, "광고 수익을 본사(40%), 가맹 운영자(40%), 입점 장소제공자(20%)가 합리적으로 분배하여, 기기의 설치 확장이 신속하고 자발적으로 이뤄지는 유통 생태계를 확보합니다.", 13, COLOR_DARK_TEXT)

    # 우측: 3대 수익 모델 구조 다이어그램
    right_left = Inches(7.2)
    models = [
        {"title": "① 장비 판매수익 (1차 수익)", "desc": "자판기 유통 마진 (대당 150만 원 이익)"},
        {"title": "② 라면 판매수익 (반복 수익)", "desc": "식재료 및 자원 공급 마진 (개당 2,350원 이익)"},
        {"title": "③ LCD 광고수익 (플랫폼 수익)", "desc": "대기시간 광고 송출 및 수수료 수취"}
    ]
    
    for i, model in enumerate(models):
        top = Inches(1.8) + i * Inches(1.6)
        
        # 다이어그램 박스 (네이비 & 레드 포인트)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_left, top, Inches(5.6), Inches(1.3))
        style_shape(box, COLOR_WHITE, COLOR_LIGHT_BG)
        
        # 왼쪽 세로 바 (레드)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_left, top, Inches(0.15), Inches(1.3))
        style_shape(bar, COLOR_PRIMARY_RED if i == 2 else COLOR_SECONDARY_NAVY)
        
        tf_model = slide.shapes.add_textbox(right_left + Inches(0.3), top + Inches(0.1), Inches(5.1), Inches(1.1)).text_frame
        tf_model.word_wrap = True
        
        p_mt = tf_model.paragraphs[0]
        set_text_formatting(p_mt, model["title"], 16, COLOR_PRIMARY_RED if i == 2 else COLOR_SECONDARY_NAVY, bold=True, font_name=FONT_TITLE)
        
        p_md = tf_model.add_paragraph()
        set_text_formatting(p_md, model["desc"], 12, COLOR_DARK_TEXT)


def build_slide_05_product_cms(prs):
    """Slide 5: Product & CMS (제품 및 광고 편성)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_common_slide_decorations(prs, slide, "02. PRODUCT & SERVICE", "16:9 LCD 탑재 라면 자판기 및 6구간 편성 전략")

    # 좌측: 기기 안내
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.5), Inches(4.8))
    ltf = left_box.text_frame
    ltf.word_wrap = True
    p0 = ltf.paragraphs[0]
    set_text_formatting(p0, "📺 하드웨어 스펙 및 특장점", 18, COLOR_SECONDARY_NAVY, bold=True, font_name=FONT_TITLE)
    p0_desc = ltf.add_paragraph()
    set_text_formatting(p0_desc, "\n• 16:9 고해상도 LCD 디스플레이 탑재\n• 실시간 광고 송출 제어를 위한 임베디드 모듈\n• 조리 진행률 안내 및 안전 수칙 실시간 표기\n• 광고 시청 후 스캔 가능한 QR 쿠폰/이벤트 연동\n• 1회 조리 시 4~8개 광고 유연 편성 지원", 13, COLOR_DARK_TEXT)

    # CMS 이미지/도식 영역 (가운데 및 오른쪽)
    # 우측: 6구간 광고 편성 프로세스 시각화
    right_left = Inches(5.3)
    p_title = slide.shapes.add_textbox(right_left, Inches(1.8), Inches(7.5), Inches(0.5))
    set_text_formatting(p_title.text_frame.paragraphs[0], "📋 조리 3분 동안의 체계적 LCD 화면 구성 (6구간)", 16, COLOR_PRIMARY_RED, bold=True, font_name=FONT_TITLE)
    
    stages = [
        {"num": "1구간", "title": "조리 시작 안내", "dur": "10s", "desc": "안전한 이용 방법 및\n대기 준비"},
        {"num": "2구간", "title": "광고 2~3개 송출", "dur": "45s", "desc": "15초 타겟 광고\n순차적 재생"},
        {"num": "3구간", "title": "조리/안전 안내", "dur": "15s", "desc": "진행률 피드백 및\n화상 주의 경고"},
        {"num": "4구간", "title": "광고 2~3개 송출", "dur": "45s", "desc": "2차 타겟 광고\n집중 노출"},
        {"num": "5구간", "title": "QR 쿠폰 & 이벤트", "dur": "20s", "desc": "브랜드 쿠폰 노출 및\n스마트폰 스캔 유도"},
        {"num": "6구간", "title": "조리 완료 안내", "dur": "15s", "desc": "음식 취출 유도 및\n감사 인사"}
    ]
    
    # 2행 3열로 6구간 시각화
    for i, stage in enumerate(stages):
        row = i // 3
        col = i % 3
        
        c_left = right_left + col * Inches(2.5)
        c_top = Inches(2.6) + row * Inches(2.1)
        
        # 개별 단계 박스
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, c_left, c_top, Inches(2.3), Inches(1.9))
        style_shape(box, COLOR_WHITE, COLOR_LIGHT_BG)
        
        # 상단 타이틀 바
        t_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, c_left, c_top, Inches(2.3), Inches(0.35))
        style_shape(t_bar, COLOR_SECONDARY_NAVY if i not in [1, 3] else COLOR_PRIMARY_RED)
        
        # 구간 텍스트
        p_sec = t_bar.text_frame.paragraphs[0]
        set_text_formatting(p_sec, f"{stage['num']} ({stage['dur']})", 11, COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        
        # 상세 설명 텍스트
        tf_desc = slide.shapes.add_textbox(c_left + Inches(0.1), c_top + Inches(0.4), Inches(2.1), Inches(1.4)).text_frame
        tf_desc.word_wrap = True
        
        p_t = tf_desc.paragraphs[0]
        set_text_formatting(p_t, stage["title"], 13, COLOR_SECONDARY_NAVY, bold=True, font_name=FONT_TITLE, alignment=PP_ALIGN.CENTER)
        
        p_d = tf_desc.add_paragraph()
        set_text_formatting(p_d, "\n" + stage["desc"], 11, COLOR_DARK_TEXT, alignment=PP_ALIGN.CENTER)


def build_slide_06_business_model(prs):
    """Slide 6: Business Model (비즈니스 모델)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_common_slide_decorations(prs, slide, "03. BUSINESS MODEL", "1대당 거래마진 및 광고 플랫폼 상생 배분 구조")

    # 3대 핵심 비즈니스 축 (3열 카드 레이아웃)
    models = [
        {
            "title": "자판기 장비 판매",
            "stat": "마진율 30.0%",
            "details": [
                "매입 원가: 350만 원",
                "판매 가격: 500만 원",
                "대당 유통이익: 150만 원",
                "매입가 기준 마진율: 42.9%",
                "초기 하드웨어 보급을 통한 빠른 본사 매출 확보 가능"
            ],
            "highlight": "150만 원 이익"
        },
        {
            "title": "라면 식자재 공급",
            "stat": "마진율 58.7%",
            "details": [
                "매입 원가: 1,650원",
                "판매 가격: 4,000원",
                "개당 유통이익: 2,350원",
                "매입가 대비 마진율: 142.4%",
                "기기 가동에 따른 고정적 원재료 반복 매출 발생 (원재료 공급 독점)"
            ],
            "highlight": "2,350원 이익"
        },
        {
            "title": "광고 매출 상생 배분",
            "stat": "본사 지분 40%",
            "details": [
                "본사(플랫폼 운영): 40%",
                "자판기 운영자: 40%",
                "장소 제공자: 20%",
                "합리적 분배로 신규 설치 부지 확보 용이",
                "CMS 솔루션 비용 및 데이터 분석 서비스 추가 창출"
            ],
            "highlight": "본사 40% 수취"
        }
    ]

    for i, bm in enumerate(models):
        left = Inches(0.5) + i * Inches(4.2)
        top = Inches(1.8)
        width = Inches(3.9)
        height = Inches(4.8)

        # 카드
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        style_shape(card, COLOR_WHITE, COLOR_LIGHT_BG)

        # 카드 헤더 (네이비)
        c_header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.7))
        style_shape(c_header, COLOR_SECONDARY_NAVY)
        p_ch = c_header.text_frame.paragraphs[0]
        set_text_formatting(p_ch, bm["title"], 18, COLOR_WHITE, bold=True, font_name=FONT_TITLE, alignment=PP_ALIGN.CENTER)

        # 내용
        tf = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.8), width - Inches(0.4), height - Inches(0.9)).text_frame
        tf.word_wrap = True

        p_stat = tf.paragraphs[0]
        set_text_formatting(p_stat, bm["stat"], 20, COLOR_PRIMARY_RED, bold=True, alignment=PP_ALIGN.CENTER)

        p_space = tf.add_paragraph()
        p_space.font.size = Pt(6)

        for detail in bm["details"]:
            p_det = tf.add_paragraph()
            set_text_formatting(p_det, "• " + detail, 12, COLOR_DARK_TEXT)

        # 하단 핵심 요약 강조 상자
        highlight_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.3), top + Inches(4.1), width - Inches(0.6), Inches(0.5))
        style_shape(highlight_box, COLOR_PRIMARY_RED)
        p_hl = highlight_box.text_frame.paragraphs[0]
        set_text_formatting(p_hl, bm["highlight"], 14, COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)


def build_slide_07_revenue_sim_1(prs):
    """Slide 7: Revenue Simulation 1 (1대 기준 월 수익 시뮬레이션)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_common_slide_decorations(prs, slide, "03. REVENUE SIMULATION", "자판기 1대당 일 판매량 기준 회수기간 및 결합 수익 분석")

    # 상단 설명 텍스트박스
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.33), Inches(0.6))
    set_text_formatting(desc_box.text_frame.paragraphs[0], "💡 자판기를 500만 원에 구매한 운영자가 라면 판매마진만으로 장비값을 전액 회수하는 기간 및 광고 결합 월 수익 시뮬레이션", 14, COLOR_PRIMARY_RED, bold=True)

    # 데이터 테이블 생성 (7행 5열)
    rows = 6
    cols = 5
    left = Inches(0.5)
    top = Inches(2.0)
    width = Inches(12.33)
    height = Inches(3.4)

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # 열 너비 세팅
    table.columns[0].width = Inches(2.0) # 하루 판매량
    table.columns[1].width = Inches(2.5) # 일 매출총이익
    table.columns[2].width = Inches(2.8) # 월 매출총이익 (30일)
    table.columns[3].width = Inches(2.5) # 장비값 회수 기간
    table.columns[4].width = Inches(2.53) # 광고 결합 월 합산 수익

    # 헤더 텍스트 설정
    headers = ["하루 판매량", "일 매출총이익", "월 매출총이익 (30일)", "순수 장비값 회수기간", "광고 월 20만 원 결합 수익"]
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_SECONDARY_NAVY
        p = cell.text_frame.paragraphs[0]
        set_text_formatting(p, header, 13, COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 로우 데이터
    row_data = [
        ["20개", "47,000원", "141.0만 원", "약 106일", "161.0만 원"],
        ["30개", "70,500원", "211.5만 원", "약 71일", "231.5만 원"],
        ["50개 (기본 목표)", "117,500원", "352.5만 원", "약 43일", "372.5만 원"],
        ["80개", "188,000원", "564.0만 원", "약 27일", "584.0만 원"],
        ["100개", "235,000원", "705.0만 원", "약 21일", "725.0만 원"]
    ]

    for row_idx, row_vals in enumerate(row_data):
        for col_idx, val in enumerate(row_vals):
            cell = table.cell(row_idx + 1, col_idx)
            cell.fill.solid()
            # 50개 기본 목표 강조
            if row_idx == 2:
                cell.fill.fore_color.rgb = RGBColor(255, 230, 230)
                font_color = COLOR_PRIMARY_RED
                is_bold = True
            else:
                cell.fill.fore_color.rgb = COLOR_WHITE if row_idx % 2 == 0 else COLOR_LIGHT_BG
                font_color = COLOR_DARK_TEXT
                is_bold = False
            
            p = cell.text_frame.paragraphs[0]
            set_text_formatting(p, val, 12, font_color, bold=is_bold, alignment=PP_ALIGN.CENTER)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 하단 유의 사항
    note_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.6), Inches(12.33), Inches(1.2))
    ntf = note_box.text_frame
    ntf.word_wrap = True
    set_text_formatting(ntf.paragraphs[0], "📌 실제 현장 참고 사항", 14, COLOR_SECONDARY_NAVY, bold=True)
    p_n1 = ntf.add_paragraph()
    set_text_formatting(p_n1, "• 실제 가맹점 운영 시에는 설치 위치에 따른 장소 수수료, 전기료, 통신비, A/S 유지 보수비 등이 차감되므로 실 회수 기간은 기재된 수치보다 연장될 수 있습니다.", 11, COLOR_DARK_TEXT)
    p_n2 = ntf.add_paragraph()
    set_text_formatting(p_n2, "• 하루 50개 판매 달성 시 광고 결합으로 대당 월 370만 원을 초과하는 합산 총이익이 발생하며 가맹점 및 본사의 견고한 현금 파이프라인이 보장됩니다.", 11, COLOR_DARK_TEXT)


def build_slide_08_revenue_sim_2(prs):
    """Slide 8: Revenue Simulation 2 (100대 보급 시 본사 수익 시뮬레이션)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_common_slide_decorations(prs, slide, "03. REVENUE SIMULATION", "100대 구축 시 본사 플랫폼 매출 및 시너지 규모")

    # 좌측: 100대 장비 판매 및 라면 유통 수익
    left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.8), Inches(5.8), Inches(4.8))
    style_shape(left_card, COLOR_WHITE, COLOR_LIGHT_BG)
    
    # 붉은색 타이틀
    t_left = slide.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(5.4), Inches(4.3))
    ltf = t_left.text_frame
    ltf.word_wrap = True
    set_text_formatting(ltf.paragraphs[0], "📈 장비 유통 및 원재료 독점 공급 매출", 16, COLOR_PRIMARY_RED, bold=True, font_name=FONT_TITLE)
    
    p1 = ltf.add_paragraph()
    set_text_formatting(p1, "\n1. 장비 판매 즉시 총이익: 1억 5천만 원", 14, COLOR_SECONDARY_NAVY, bold=True)
    p1_sub = ltf.add_paragraph()
    set_text_formatting(p1_sub, "   - 100대 판매 시 매출 5억 원, 매입원가 3.5억 원\n   - 초기 대규모 현금 흐름 확보로 인증/솔루션 R&D 자금 충당", 12, COLOR_DARK_TEXT)
    
    p2 = ltf.add_paragraph()
    set_text_formatting(p2, "\n2. 라면 식재료 월 매출총이익: 3.52억 원 (전체 기준)", 14, COLOR_SECONDARY_NAVY, bold=True)
    p2_sub = ltf.add_paragraph()
    set_text_formatting(p2_sub, "   - 하루 50개 판매 시 100대 총 월 150,000개 유통\n   - 월 라면 매출 6억 원, 매출총이익 3.52억 원 규모\n   - 본사는 라면 전용 스프 및 기기용 포장 식자재 공급 마진 수취", 12, COLOR_DARK_TEXT)

    # 우측: 100대 보급 시 광고수익 테이블 시뮬레이션
    right_left = Inches(6.8)
    r_title = slide.shapes.add_textbox(right_left, Inches(1.8), Inches(6.0), Inches(0.5))
    set_text_formatting(r_title.text_frame.paragraphs[0], "📺 100대 기준 광고 매출 및 본사 몫 (40%)", 16, COLOR_SECONDARY_NAVY, bold=True, font_name=FONT_TITLE)

    rows = 5
    cols = 3
    table_shape = slide.shapes.add_table(rows, cols, right_left, Inches(2.4), Inches(6.0), Inches(3.6))
    table = table_shape.table
    table.columns[0].width = Inches(2.0) # 대당 월 광고료
    table.columns[1].width = Inches(2.0) # 100대 월 총광고매출
    table.columns[2].width = Inches(2.0) # 본사 40% 플랫폼수익

    # 테이블 헤더
    t_headers = ["대당 월 광고료", "100대 총 매출", "본사 수수료 (40%)"]
    for col_idx, header in enumerate(t_headers):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_SECONDARY_NAVY
        set_text_formatting(cell.text_frame.paragraphs[0], header, 12, COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    table_data = [
        ["10만 원", "1,000만 원", "월 400만 원"],
        ["20만 원 (목표)", "2,000만 원", "월 800만 원"],
        ["30만 원", "3,000만 원", "월 1,200만 원"],
        ["50만 원 (핫스팟)", "5,000만 원", "월 2,000만 원"]
    ]

    for row_idx, row_vals in enumerate(table_data):
        for col_idx, val in enumerate(row_vals):
            cell = table.cell(row_idx + 1, col_idx)
            cell.fill.solid()
            if row_idx == 1:
                cell.fill.fore_color.rgb = RGBColor(255, 230, 230)
                font_color = COLOR_PRIMARY_RED
                is_bold = True
            else:
                cell.fill.fore_color.rgb = COLOR_WHITE if row_idx % 2 == 0 else COLOR_LIGHT_BG
                font_color = COLOR_DARK_TEXT
                is_bold = False
            set_text_formatting(cell.text_frame.paragraphs[0], val, 11, font_color, bold=is_bold, alignment=PP_ALIGN.CENTER)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def build_slide_09_market_opportunity(prs):
    """Slide 9: Market Opportunity (시장 기회 및 타겟 스페이스)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_common_slide_decorations(prs, slide, "04. MARKET OPPORTUNITY", "핵심 설치 후보 공간 및 광고주 타겟 시장 분석")

    # 좌측: 핵심 입점 타겟 스페이스 (5대 최적 설치 장소)
    left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.8), Inches(5.8), Inches(4.8))
    style_shape(left_card, COLOR_WHITE, COLOR_LIGHT_BG)
    
    tf_l = slide.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(5.4), Inches(4.4)).text_frame
    tf_l.word_wrap = True
    set_text_formatting(tf_l.paragraphs[0], "📍 5대 핵심 밀착형 입점 공간", 18, COLOR_PRIMARY_RED, bold=True, font_name=FONT_TITLE)
    
    spaces = [
        ("대학교 기숙사", "야간/새벽 시간 야식 수요 집중 및 대학생 대상 캠퍼스 광고 최적지"),
        ("PC방 / 멀티방", "게임 대기 및 이용 중 간편 조리식 수요 연동, 게임/IT 브랜드 연계 광고"),
        ("제조 공장 / 휴게공간", "상시 근무하는 블루칼라 직원의 야식 보장 및 안정적 회전율 보장"),
        ("고시원 / 스터디카페", "장시간 체류하는 학생, 수험생의 간편 야식 수요 및 교육 콘텐츠 광고"),
        ("병원 / 요양시설", "교대 근무 간호사 및 상주 직원의 심야 식사 확보 및 건강식 위주 타겟팅")
    ]
    
    for title, desc in spaces:
        p = tf_l.add_paragraph()
        set_text_formatting(p, f"\n✔ {title}", 13, COLOR_SECONDARY_NAVY, bold=True)
        p_sub = tf_l.add_paragraph()
        set_text_formatting(p_sub, f"   {desc}", 11, COLOR_DARK_TEXT)

    # 우측: 타겟 광고주 및 패키지 구조
    right_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(6.0), Inches(4.8))
    style_shape(right_card, COLOR_WHITE, COLOR_LIGHT_BG)

    tf_r = slide.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.6), Inches(4.4)).text_frame
    tf_r.word_wrap = True
    set_text_formatting(tf_r.paragraphs[0], "🎯 로컬 소상공인 & 프랜차이즈 브랜드 결합", 18, COLOR_SECONDARY_NAVY, bold=True, font_name=FONT_TITLE)

    ads = [
        ("지역 소상공인 밀착형 광고", "자판기 주변 반경 500m 이내의 요식업, 뷰티, 피트니스 등 소상공인이 모바일을 통해 손쉽게 쿠폰 스캔 형태로 홍보할 수 있는 '로컬 베이직' 상품 출시"),
        ("프랜차이즈 및 브랜드 쿠폰 연동", "광고 시청 완료 후 화면상의 QR코드를 스캔하면 전국 프랜차이즈 매장에서 바로 사용 가능한 즉시 할인 쿠폰을 지급하는 프로모션형 브랜드 패키지 광고 유치"),
        ("공공 및 지역 안전 캠페인", "정부, 구청 등 공공기관의 청년/지역 일자리 정책 및 소방/보건안전 캠페인을 자판기 조리 대기 시간 동안 의무 송출하여 사회적 신뢰 매체로 포지셔닝")
    ]

    for title, desc in ads:
        p = tf_r.add_paragraph()
        set_text_formatting(p, f"\n📢 {title}", 13, COLOR_PRIMARY_RED, bold=True)
        p_sub = tf_r.add_paragraph()
        set_text_formatting(p_sub, f"   {desc}", 11, COLOR_DARK_TEXT)


def build_slide_10_gtm_strategy(prs):
    """Slide 10: Go-To-Market Strategy (단계별 성장 로드맵)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_common_slide_decorations(prs, slide, "04. GO-TO-MARKET", "10대 파일럿에서 전국 미디어 플랫폼으로의 5단계 성장 로드맵")

    # 5단계 가로형 화살표/박스 타임라인 시각화
    steps = [
        {"step": "1단계", "title": "제품 인증 및 규격화", "desc": "• KC안전/전파 인증\n• 식품위생 자판기 신고\n• CMS 기본 R&D 완료"},
        {"step": "2단계", "title": "10대 파일럿 검증", "desc": "• 기숙사, PC방 등 설치\n• 판매수량/노출 데이터 확보\n• 이용자 QR 스캔율 분석"},
        {"step": "3단계", "title": "지역 광고 상품 출시", "desc": "• 광고 CMS 고도화\n• 소상공인 연동 상품 개발\n• 30대 규모로 보급 확장"},
        {"step": "4단계", "title": "100대 지역 거점 구축", "desc": "• 대리점/위탁운영 체계\n• 플랫폼 수수료 계약 체결\n• 지역 밀착형 영업망 완성"},
        {"step": "5단계", "title": "전국 미디어 플랫폼", "desc": "• 500대 이상 전국망 확보\n• 프랜차이즈 대기업 광고\n• 빅데이터 DOOH 플랫폼화"}
    ]

    width_step = Inches(2.2)
    height_step = Inches(4.3)
    gap = Inches(0.2)

    for i, step in enumerate(steps):
        left = Inches(0.5) + i * (width_step + gap)
        top = Inches(1.8)

        # 개별 단계 사각형 카드
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width_step, height_step)
        style_shape(card, COLOR_WHITE, COLOR_LIGHT_BG)

        # 상단 단계 표시 바 (레드 & 네이비 교차 또는 1단계 강조)
        t_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width_step, Inches(0.6))
        style_shape(t_bar, COLOR_PRIMARY_RED if i < 2 else COLOR_SECONDARY_NAVY)
        p_t = t_bar.text_frame.paragraphs[0]
        set_text_formatting(p_t, step["step"], 14, COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        # 내용 작성
        tf = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.7), width_step - Inches(0.2), height_step - Inches(0.8)).text_frame
        tf.word_wrap = True

        p_title = tf.paragraphs[0]
        set_text_formatting(p_title, step["title"], 13, COLOR_SECONDARY_NAVY, bold=True, font_name=FONT_TITLE, alignment=PP_ALIGN.CENTER)

        p_desc = tf.add_paragraph()
        set_text_formatting(p_desc, "\n" + step["desc"], 11, COLOR_DARK_TEXT)

        # 단계 연결 화살표 (마지막 단계 제외)
        if i < 4:
            arrow_left = left + width_step + Inches(0.02)
            arrow_top = top + Inches(1.8)
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_left, arrow_top, Inches(0.16), Inches(0.3))
            style_shape(arrow, COLOR_PRIMARY_RED if i < 1 else COLOR_MUTED_TEXT)


def build_slide_11_investment_funds(prs):
    """Slide 11: Competitive Advantage & Investment Point (투자 가치 및 자금 사용 계획)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_common_slide_decorations(prs, slide, "04. INVESTMENT POINT", "핵심 투자 메리트 및 조달 자금 집행 계획")

    # 좌측: 핵심 경쟁 우위 및 투자 포인트
    left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.8), Inches(5.8), Inches(4.8))
    style_shape(left_card, COLOR_WHITE, COLOR_LIGHT_BG)

    tf_l = slide.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(5.4), Inches(4.4)).text_frame
    tf_l.word_wrap = True
    set_text_formatting(tf_l.paragraphs[0], "💎 핵심 투자 포인트 (Investment Point)", 18, COLOR_PRIMARY_RED, bold=True, font_name=FONT_TITLE)

    pts = [
        ("빠른 현금 흐름", "초기 자판기 장비 대당 150만 원 유통 마진 확보로 BEP 조기 달성 가능"),
        ("누적형 플랫폼 반복 수익", "기기가 설치될수록 자원 유통 마진과 광고 매출 수수료가 누적 스노우볼링"),
        ("독점적 공간 점유 효과", "입점 부지 선점 시 경쟁사 진입이 어려운 '식품 판매 + 로컬 미디어' 결합망 완성"),
        ("데이터 비즈니스 확장", "이용자 소비 행동 데이터(QR 스캔, 결제 내역) 결합형 마케팅 솔루션으로 고도화")
    ]

    for title, desc in pts:
        p = tf_l.add_paragraph()
        set_text_formatting(p, f"\n💡 {title}", 13, COLOR_SECONDARY_NAVY, bold=True)
        p_sub = tf_l.add_paragraph()
        set_text_formatting(p_sub, f"   {desc}", 11, COLOR_DARK_TEXT)

    # 우측: 자금 사용 계획 (Use of Funds)
    right_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(6.0), Inches(4.8))
    style_shape(right_card, COLOR_WHITE, COLOR_LIGHT_BG)

    tf_r = slide.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.6), Inches(4.4)).text_frame
    tf_r.word_wrap = True
    set_text_formatting(tf_r.paragraphs[0], "💰 자금 사용 계획 (Use of Funds)", 18, COLOR_SECONDARY_NAVY, bold=True, font_name=FONT_TITLE)

    funds = [
        ("기기 인증 및 위생 안전 규격 확보", "KC안전 인증, 전파 적합성 평가, 식약처 위생자동판매기 안전 기준 통과"),
        ("초기 파일럿 기기 인프라 조달", "대학가 기숙사 및 대형 PC방 설치용 10대~30대 파일럿 장비 수입 및 조립 비용"),
        ("광고 송출 CMS 솔루션 고도화", "원격 실시간 편성 조절, 지역별 매칭, 광고주의 쿠폰 발행 및 대시보드 R&D"),
        ("전국단위 영업망 및 A/S망 체계화", "기기 오작동에 대한 전국 24h 대응 체계 구축 및 직영/대리점 영업 인프라 세팅")
    ]

    for title, desc in funds:
        p = tf_r.add_paragraph()
        set_text_formatting(p, f"\n📌 {title}", 13, COLOR_PRIMARY_RED, bold=True)
        p_sub = tf_r.add_paragraph()
        set_text_formatting(p_sub, f"   {desc}", 11, COLOR_DARK_TEXT)


def build_slide_12_closing(prs):
    """Slide 12: Closing (결언)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 어두운 남색 풀배경 설정 (결언 임팩트 극대화)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    style_shape(bg, COLOR_SECONDARY_NAVY)

    # 붉은색 기하학 삼각형 데코레이션
    tri1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(9.333), Inches(3.5), Inches(4.0), Inches(4.0))
    style_shape(tri1, COLOR_PRIMARY_RED)
    tri1.rotation = 270

    # 텍스트 박스 (중앙)
    text_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.33), Inches(4.5))
    tf = text_box.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    set_text_formatting(p0, "THE VISION", 16, COLOR_PRIMARY_RED, bold=True, alignment=PP_ALIGN.CENTER)
    
    p_space1 = tf.add_paragraph()
    p_space1.font.size = Pt(20)

    p1 = tf.add_paragraph()
    set_text_formatting(p1, "“라면 자판기를 유통하는 제조사가 아닙니다.”", 26, COLOR_WHITE, bold=True, font_name=FONT_TITLE, alignment=PP_ALIGN.CENTER)

    p2 = tf.add_paragraph()
    set_text_formatting(p2, "“조리 대기시간 3분을 혁신적 미디어 공간으로 전환하는\n독점적 생활밀착형 광고 플랫폼입니다.”", 30, COLOR_PRIMARY_RED, bold=True, font_name=FONT_TITLE, alignment=PP_ALIGN.CENTER)

    p_space2 = tf.add_paragraph()
    p_space2.font.size = Pt(28)

    p3 = tf.add_paragraph()
    set_text_formatting(p3, "주식회사 엔에스애드홉 (Noodle Ads Hub) 드림", 14, COLOR_MUTED_TEXT, bold=False, alignment=PP_ALIGN.CENTER)


# ==========================================
# 메인 빌드 파이프라인
# ==========================================
def main():
    prs = Presentation()
    # 16:9 와이드스크린 슬라이드 화면 비율 설정
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    print("Building IR Slide 1 (Cover)...")
    build_slide_01_cover(prs)
    
    print("Building IR Slide 2 (TOC)...")
    build_slide_02_toc(prs)
    
    print("Building IR Slide 3 (Problem)...")
    build_slide_03_problem(prs)
    
    print("Building IR Slide 4 (Solution)...")
    build_slide_04_solution(prs)
    
    print("Building IR Slide 5 (Product & CMS)...")
    build_slide_05_product_cms(prs)
    
    print("Building IR Slide 6 (Business Model)...")
    build_slide_06_business_model(prs)
    
    print("Building IR Slide 7 (Revenue Sim 1)...")
    build_slide_07_revenue_sim_1(prs)
    
    print("Building IR Slide 8 (Revenue Sim 2)...")
    build_slide_08_revenue_sim_2(prs)
    
    print("Building IR Slide 9 (Market Opportunity)...")
    build_slide_09_market_opportunity(prs)
    
    print("Building IR Slide 10 (GTM Strategy)...")
    build_slide_10_gtm_strategy(prs)
    
    print("Building IR Slide 11 (Investment & Funds)...")
    build_slide_11_investment_funds(prs)
    
    print("Building IR Slide 12 (Closing)...")
    build_slide_12_closing(prs)

    output_path = "Noodle_Ads_Hub_IR_Pitch_Deck.pptx"
    prs.save(output_path)
    print(f"\n[성공] PPTX 프레젠테이션 생성이 완료되었습니다: {output_path}")

if __name__ == "__main__":
    main()
